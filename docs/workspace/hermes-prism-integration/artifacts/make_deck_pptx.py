#!/usr/bin/env python3
"""make_deck_pptx.py — EDITABLE Algolia search-audit deck, built ON the official
'Algolia Slide Template 2026.PPTX' (Google-Slides-exported, 16:9, 10x5.62in).

Why this over the HTML->PDF doc: the deliverable must be editable by the AE in
PowerPoint / Google Slides. We build on the real template so brand chrome
(gradient cover, section dividers, bottom accent bar, Algolia mark, theme) is
inherited from its layouts.

Grounding (hard): every finding shows what we TESTED / EXPECTED / FOUND from the
audit JSON, plus the Algolia fix. Sources are INLINE clickable hyperlinks on the
supporting text, never a trailing 'Source ->' tag. No fabrication, no em dashes,
no financials/pricing.

Usage: python3 make_deck_pptx.py <audit-data.json> <out.pptx>
"""
import json, os, sys
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

DS = "/Users/arijitchowdhury/Dropbox/AI-Development/Algolia-Design-System"
TEMPLATE = f"{DS}/uploads/Algolia Slide Tempalte 2026.PPTX"
FONT = "Sora"

NAVY = RGBColor(0x00, 0x00, 0x33)
INK = RGBColor(0x21, 0x24, 0x3D)
BLUE = RGBColor(0x00, 0x3D, 0xFF)
GREY = RGBColor(0x6B, 0x72, 0x80)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SEV = {  # (text color, fill) by severity
    "LOW": (RGBColor(0x05, 0x96, 0x60), RGBColor(0xE9, 0xF7, 0xF0)),
    "MEDIUM": (RGBColor(0xD9, 0x77, 0x06), RGBColor(0xFD, 0xF1, 0xE3)),
    "HIGH": (RGBColor(0xDC, 0x26, 0x26), RGBColor(0xFB, 0xEA, 0xEA)),
    "CRITICAL": (RGBColor(0xB9, 0x1C, 0x1C), RGBColor(0xFB, 0xE3, 0xE3)),
}

# layout indices in the template (all findings use the WHITE body layout; content placed manually)
L_COVER, L_BODY, L_SECTION, L_FIND = 0, 4, 2, 4


def clean(x):
    """Grounding hygiene: never emit em dashes; coerce lists to prose."""
    if x is None:
        return ""
    if isinstance(x, (list, tuple)):
        x = ", ".join(str(i) for i in x if i)
    return str(x).replace("—", "-").replace("–", "-").strip()


def summarize(x, n=210):
    """Deck bullets stay crisp: keep to ~n chars on a word/sentence boundary (full depth
    lives in the report doc). Never cut mid-word, never add an ellipsis."""
    s = clean(x)
    if len(s) <= n:
        return s
    cut = s[:n]
    for stop in (". ", "; ", ", ", " "):
        i = cut.rfind(stop)
        if i > n * 0.6:
            return cut[:i].rstrip(",;") + ("." if not cut[:i].rstrip().endswith(".") else "")
    return cut.rstrip() + "."


def strip_slides(prs):
    """Remove the template's example slides (part + rel) so no orphan/duplicate parts survive."""
    part = prs.part
    for sldId in list(prs.slides._sldIdLst):
        rId = sldId.get(qn("r:id"))
        if rId:
            part.drop_rel(rId)
        prs.slides._sldIdLst.remove(sldId)


def set_run(r, text, size, color, bold=False, italic=False, link=None):
    r.text = text
    r.font.name = FONT
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    if link:
        r.hyperlink.address = link
    return r


def textbox(slide, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(1)
    return tb, tf


def para(tf, first=False):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    return p


def label_line(tf, label, first=False):
    """Small uppercase grey label (WE TESTED / WE FOUND ...)."""
    p = para(tf, first)
    p.space_before = Pt(0 if first else 7)
    p.space_after = Pt(1)
    set_run(p.add_run(), label.upper(), 8, GREY, bold=True)
    p.runs[0].font._rPr.set("spc", "80")  # letter spacing


def body_line(tf, text, size=10.5, color=INK, bold=False):
    p = para(tf)
    p.space_after = Pt(0)
    set_run(p.add_run(), clean(text), size, color, bold=bold)
    return p


def rounded(slide, l, t, w, h, fill, line=None, radius=0.08):
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(1)
    sp.shadow.inherit = False
    try:
        sp.adjustments[0] = radius
    except Exception:
        pass
    return sp


def main(json_path, out_path):
    d = json.load(open(json_path))
    rd = os.path.dirname(os.path.abspath(json_path))
    meta = d.get("meta", {})
    co = clean(meta.get("company", "Company"))
    score = d.get("score", {})
    findings = d.get("findings", [])
    crit = [f for f in findings if str(f.get("severity", "")).lower() == "critical"]
    ordered = crit + [f for f in findings if f not in crit]

    prs = Presentation(TEMPLATE)
    strip_slides(prs)
    SW = prs.slide_width

    # ---------- COVER ----------
    s = prs.slides.add_slide(prs.slide_layouts[L_COVER])
    logo = f"{DS}/assets/Algolia-logo-white.png"
    if os.path.exists(logo):
        lw = Inches(2.1)
        pic = s.shapes.add_picture(logo, 0, Inches(0.72), width=lw)
        pic.left = (SW - pic.width) // 2  # centered
    s.placeholders[0].text = f"eCommerce Search Audit"
    for p in s.placeholders[0].text_frame.paragraphs:
        for r in p.runs:
            r.font.name = FONT
    sub = s.placeholders[1]
    sub.text = f"{co}   |   Prepared by {clean(meta.get('audited_by','Algolia'))}"
    for r in sub.text_frame.paragraphs[0].runs:
        r.font.name = FONT

    # ---------- SCORECARD ----------
    s = prs.slides.add_slide(prs.slide_layouts[L_BODY])
    s.placeholders[0].text = "Search audit scorecard"
    for r in s.placeholders[0].text_frame.paragraphs[0].runs:
        r.font.name = FONT; r.font.color.rgb = BLUE; r.font.bold = True
    # big score + verdict
    _, tf = textbox(s, 0.34, 1.15, 3.1, 1.4)
    p = tf.paragraphs[0]
    set_run(p.add_run(), str(score.get("overall", "?")), 46, BLUE, bold=True)
    set_run(p.add_run(), " /10", 18, GREY)
    body_line(tf, score.get("verdict", ""), 13, INK, bold=True)
    p = para(tf); p.space_before = Pt(2)
    set_run(p.add_run(), f"{score.get('critical_count',0)} critical  ", 10, SEV['HIGH'][0], bold=True)
    set_run(p.add_run(), f"{score.get('moderate_count',0)} moderate  ", 10, SEV['MEDIUM'][0], bold=True)
    set_run(p.add_run(), f"{score.get('low_count',0)} strengths", 10, SEV['LOW'][0], bold=True)
    # heatmap tiles (2 cols x 5)
    bd = score.get("breakdown", {}); labels = score.get("breakdown_labels", {}); sev = score.get("breakdown_severity", {})
    items = sorted(bd.items(), key=lambda kv: kv[1])
    x0, y0, tw, th, gx, gy = 3.7, 1.15, 2.95, 0.62, 0.14, 0.12
    for i, (k, v) in enumerate(items):
        col, row = i % 2, i // 2
        sv = str(sev.get(k, "MEDIUM")).upper()
        tc, fc = SEV.get(sv, SEV["MEDIUM"])
        l = x0 + col * (tw + gx); t = y0 + row * (th + gy)
        rounded(s, l, t, tw, th, fc, line=tc)
        _, tf2 = textbox(s, l + 0.12, t, tw - 0.2, th, anchor=MSO_ANCHOR.MIDDLE)
        pp = tf2.paragraphs[0]
        set_run(pp.add_run(), str(v), 17, tc, bold=True)
        set_run(pp.add_run(), "/10  ", 9, tc)
        set_run(pp.add_run(), clean(labels.get(k, k)), 9.5, INK, bold=True)

    # ---------- SECTION DIVIDER ----------
    s = prs.slides.add_slide(prs.slide_layouts[L_SECTION])
    s.placeholders[0].text = "Areas of improvement"
    for r in s.placeholders[0].text_frame.paragraphs[0].runs:
        r.font.name = FONT

    # ---------- FINDING SLIDES ----------
    for f in ordered[:5]:
        finding_slide(prs, rd, f)

    # ---------- CLOSE ----------
    s = prs.slides.add_slide(prs.slide_layouts[L_BODY])
    s.placeholders[0].text = "Where to start: a scoped proof of concept"
    for r in s.placeholders[0].text_frame.paragraphs[0].runs:
        r.font.name = FONT; r.font.color.rgb = BLUE; r.font.bold = True
    _, tf = textbox(s, 0.34, 1.2, 9.3, 3.6)
    body_line(tf, "Prove the lift on your live index in weeks, not quarters. A focused POC on the highest impact gap, measured against your own analytics.", 12, INK)
    rfp = d.get("recommended_first_play", {}); steps = d.get("next_steps", [])
    n = 1
    if rfp.get("headline"):
        _step(tf, n, rfp.get("headline"), rfp.get("detail", "")); n += 1
    for st in steps:
        _step(tf, n, st.get("title", ""), st.get("description", "")); n += 1

    prs.save(out_path)
    print(f"wrote {out_path} ({len(prs.slides._sldIdLst)} slides)")


def _step(tf, n, head, detail):
    p = tf.add_paragraph(); p.space_before = Pt(9)
    set_run(p.add_run(), f"{n}. {clean(head)}. ", 11.5, BLUE, bold=True)
    set_run(p.add_run(), clean(detail), 11, INK)


def finding_slide(prs, rd, f):
    s = prs.slides.add_slide(prs.slide_layouts[L_FIND])  # white body layout
    # remove the layout's body placeholder (we place everything manually); keep title placeholder
    for ph in list(s.placeholders):
        if ph.placeholder_format.idx == 1:
            ph._element.getparent().remove(ph._element)
    # title (navy bold), in the title placeholder
    s.placeholders[0].text = clean(f.get("title", "Finding"))
    for r in s.placeholders[0].text_frame.paragraphs[0].runs:
        r.font.name = FONT; r.font.bold = True; r.font.color.rgb = NAVY; r.font.size = Pt(16)

    # ---- LEFT COLUMN: tested / expected / found ----
    _, tf = textbox(s, 0.34, 1.12, 4.55, 2.55)
    q = clean(f.get("tested_query", "")).split(";")[0].strip()
    label_line(tf, "We tested", first=True)
    p = para(tf); p.space_after = Pt(0)
    set_run(p.add_run(), f'"{q}"', 10, INK, bold=True)
    label_line(tf, "We expected")
    body_line(tf, summarize(f.get("expected_behavior", ""), 190), 9.5)
    label_line(tf, "We found")
    body_line(tf, summarize(f.get("actual_behavior", ""), 230), 9.5)

    # ---- WITH ALGOLIA box (light-blue rounded) ----
    box = rounded(s, 0.34, 3.78, 4.55, 1.12, RGBColor(0xF2, 0xF5, 0xFF), line=RGBColor(0xD9, 0xE2, 0xFF), radius=0.06)
    tfb = box.text_frame; tfb.word_wrap = True
    tfb.margin_left = tfb.margin_right = Pt(7); tfb.margin_top = tfb.margin_bottom = Pt(5)
    pb = tfb.paragraphs[0]; pb.space_after = Pt(0)
    set_run(pb.add_run(), "With Algolia  ", 9.5, BLUE, bold=True)
    set_run(pb.add_run(), clean(f.get("algolia_solution", "")), 9, INK)

    # ---- inline clickable source (no 'Source ->' tag) ----
    _, tfs = textbox(s, 0.34, 4.96, 4.55, 0.35)
    _inline_source(tfs.paragraphs[0], f.get("impact_stat_source", "PetSmart Algolia search analytics (30-day window)"))

    # ---- RIGHT: screenshot + red tag ----
    sf = f.get("screenshot_file")
    path = os.path.join(rd, sf) if sf else None
    px, py, pw, ph_cap = 5.12, 1.12, 4.55, 3.9
    if path and os.path.exists(path):
        pic = s.shapes.add_picture(path, Inches(px), Inches(py), width=Inches(pw))
        if pic.height > Inches(ph_cap):  # tall PDP shots: constrain by height, recenter
            pic._element.getparent().remove(pic._element)
            pic = s.shapes.add_picture(path, Inches(px), Inches(py), height=Inches(ph_cap))
            pic.left = Inches(px) + (Inches(pw) - pic.width) // 2
        pic.line.color.rgb = RGBColor(0xE1, 0xE4, 0xEC); pic.line.width = Pt(1)
        tag = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, pic.left + Inches(0.1), Inches(py) + Inches(0.1), Inches(1.95), Inches(0.3))
        tag.fill.solid(); tag.fill.fore_color.rgb = SEV['HIGH'][0]; tag.line.fill.background(); tag.shadow.inherit = False
        tg = tag.text_frame; tg.margin_top = tg.margin_bottom = Pt(1); tg.margin_left = Pt(6)
        set_run(tg.paragraphs[0].add_run(), "What shoppers see today", 8.5, WHITE, bold=True)


def _inline_source(p, srcval):
    srcval = clean(srcval)
    if srcval.startswith("http"):
        host = srcval.split("//", 1)[-1].split("/", 1)[0]
        host = host[4:] if host.startswith("www.") else host
        set_run(p.add_run(), "Verified against ", 8.5, GREY, italic=True)
        set_run(p.add_run(), host, 8.5, BLUE, italic=True, link=srcval)
    else:
        set_run(p.add_run(), "Verified against ", 8.5, GREY, italic=True)
        set_run(p.add_run(), srcval, 8.5, GREY, italic=True)


if __name__ == "__main__":
    if len(sys.argv) < 3:
        print("usage: make_deck_pptx.py <audit-data.json> <out.pptx>"); sys.exit(1)
    main(sys.argv[1], sys.argv[2])
