from pptx import Presentation
from pptx.util import Emu
p = Presentation("/Users/arijitchowdhury/Dropbox/AI-Development/Algolia-Design-System/uploads/Algolia Slide Tempalte 2026.PPTX")
print("slide size:", round(p.slide_width/914400,2), "x", round(p.slide_height/914400,2), "in")
print("layouts:", len(p.slide_layouts))
for i, lay in enumerate(p.slide_layouts):
    phs = []
    for ph in lay.placeholders:
        phs.append(f"[{ph.placeholder_format.idx}]{ph.placeholder_format.type} '{ph.name}'")
    print(f"\n#{i} '{lay.name}'")
    if phs:
        for x in phs: print("   ", x)
    else:
        print("    (no placeholders)")
